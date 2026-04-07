"""
Fly OS — Triple Hybrid Engine Client v5.0

Engine 1 (Primary):   Google Gemini 2.5 Flash  — multimodal, streaming, fast
Engine 2 (Secondary): OpenRouter Qwen3          — fallback
Engine 3 (Tertiary):  Base44 Agent              — final fallback

New in v5:
  - URL/webpage reading (auto-detect pasted links)
  - Video frame extraction and analysis
  - PPTX / XLSX generation
  - Stop/cancel mid-stream
  - Response mode: fast / balanced / research / creative
  - Web search grounding (Gemini native)
  - Background Python code execution
  - Auto-update version check
"""

import threading
import sqlite3
import os
import base64
import mimetypes
import requests
import time
import re
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Callable, Optional
from datetime import datetime, date

TODAY = date.today().strftime("%A, %B %d, %Y")

from google import genai
from google.genai import types

from app.config import (
    GEMINI_API_KEY, OPENROUTER_API_KEY, BASE44_API_KEY,
    BASE44_AGENT_ID, BASE44_CONVO_ID,
    MODEL_GEMINI, MODEL_OPENROUTER,
    MAX_TOKENS, TEMPERATURE,
    DB_PATH, WORKSPACE_DIR, MODE_PROMPTS, RESPONSE_MODES,
    IMAGE_EXTENSIONS, DOCUMENT_EXTENSIONS, VIDEO_EXTENSIONS, AUDIO_EXTENSIONS,
    APP_VERSION, GITHUB_VERSION_URL,
)


# ── URL detection ──────────────────────────────────────────────────────────────

URL_RE = re.compile(
    r'https?://[^\s<>"\']+|www\.[^\s<>"\']+',
    re.IGNORECASE
)

def extract_urls(text: str) -> list[str]:
    urls = URL_RE.findall(text)
    result = []
    for u in urls:
        if not u.startswith("http"):
            u = "https://" + u
        result.append(u)
    return result

def fetch_url_content(url: str, timeout: int = 15) -> str:
    """Fetch and extract text content from a URL."""
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                          "AppleWebKit/537.36 (KHTML, like Gecko) "
                          "Chrome/122.0 Safari/537.36",
            "Accept": "text/html,application/xhtml+xml,*/*;q=0.9",
        }
        resp = requests.get(url, headers=headers, timeout=timeout)
        resp.raise_for_status()
        content_type = resp.headers.get("Content-Type", "")
        if "text/html" in content_type:
            try:
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(resp.text, "html.parser")
                # Remove scripts, styles, nav
                for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
                    tag.decompose()
                text = soup.get_text(separator="\n", strip=True)
                # Collapse excessive blank lines
                lines = [l for l in text.splitlines() if l.strip()]
                return "\n".join(lines[:800])  # limit to ~800 lines
            except ImportError:
                return resp.text[:8000]
        elif "application/json" in content_type:
            return resp.text[:8000]
        elif "text/" in content_type:
            return resp.text[:8000]
        else:
            return f"[Binary content at {url} — type: {content_type}]"
    except Exception as e:
        return f"[Could not fetch {url}: {e}]"


# ── Video frame extraction ─────────────────────────────────────────────────────

def extract_video_frames(video_path: str, max_frames: int = 8) -> list[str]:
    """Extract evenly-spaced frames from a video as PNG paths."""
    try:
        import cv2
        cap = cv2.VideoCapture(video_path)
        total = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        fps   = cap.get(cv2.CAP_PROP_FPS) or 24
        duration = total / fps

        frame_indices = [int(i * total / max_frames) for i in range(max_frames)]
        paths = []
        for idx in frame_indices:
            cap.set(cv2.CAP_PROP_POS_FRAMES, idx)
            ret, frame = cap.read()
            if ret:
                ts = int(idx / fps)
                fname = f"frame_{ts:04d}s.png"
                fpath = os.path.join(WORKSPACE_DIR, fname)
                cv2.imwrite(fpath, frame)
                paths.append(fpath)
        cap.release()
        return paths
    except ImportError:
        return []
    except Exception:
        return []


# ── Document export ────────────────────────────────────────────────────────────

def export_docx(text: str, filename: str) -> str:
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

        doc = Document()

        # Page margins
        section = doc.sections[0]
        section.top_margin    = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin   = Inches(1.2)
        section.right_margin  = Inches(1.2)

        # Styles
        normal = doc.styles["Normal"]
        normal.font.name = "Calibri"
        normal.font.size = Pt(11)

        # Title
        title = doc.add_heading("Fly OS Export", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.LEFT
        title.runs[0].font.color.rgb = RGBColor(0x6C, 0x63, 0xFF)

        meta = doc.add_paragraph(f"Generated by Fly OS v5  ·  {TODAY}")
        meta.runs[0].font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
        meta.runs[0].font.size = Pt(9)

        # Horizontal rule (table trick)
        tbl = doc.add_table(rows=1, cols=1)
        tbl.style = "Table Grid"
        tbl.rows[0].cells[0].text = ""
        tbl.rows[0].height = Pt(1)
        doc.add_paragraph("")

        for line in text.split("\n"):
            if line.startswith("# "):
                h = doc.add_heading(line[2:], level=1)
                h.runs[0].font.color.rgb = RGBColor(0x6C, 0x63, 0xFF)
            elif line.startswith("## "):
                doc.add_heading(line[3:], level=2)
            elif line.startswith("### "):
                doc.add_heading(line[4:], level=3)
            elif line.startswith("```"):
                p = doc.add_paragraph(line)
                p.runs[0].font.name = "Courier New"
                p.runs[0].font.size = Pt(9)
            else:
                doc.add_paragraph(line)

        path = os.path.join(WORKSPACE_DIR, filename)
        doc.save(path)
        return path
    except ImportError:
        raise RuntimeError("python-docx not installed.")


def export_pdf(text: str, filename: str) -> str:
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable, Preformatted
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.colors import HexColor
        from reportlab.lib.enums import TA_LEFT

        path = os.path.join(WORKSPACE_DIR, filename)
        doc = SimpleDocTemplate(path, pagesize=letter,
                                leftMargin=inch*1.2, rightMargin=inch*1.2,
                                topMargin=inch, bottomMargin=inch)
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle("FlyTitle", parent=styles["Heading1"],
                                     textColor=HexColor("#6C63FF"), fontSize=20,
                                     spaceAfter=4, leading=26)
        meta_style  = ParagraphStyle("FlyMeta", parent=styles["Normal"],
                                     textColor=HexColor("#64748B"), fontSize=9, spaceAfter=8)
        h1_style    = ParagraphStyle("FlyH1",   parent=styles["Heading1"],
                                     textColor=HexColor("#6C63FF"), fontSize=16, spaceAfter=6)
        h2_style    = ParagraphStyle("FlyH2",   parent=styles["Heading2"],
                                     fontSize=13, spaceAfter=4)
        h3_style    = ParagraphStyle("FlyH3",   parent=styles["Heading3"],
                                     fontSize=11, spaceAfter=4)
        body_style  = ParagraphStyle("FlyBody", parent=styles["Normal"],
                                     fontSize=11, leading=17, spaceAfter=4)
        code_style  = ParagraphStyle("FlyCode", parent=styles["Normal"],
                                     fontName="Courier", fontSize=9,
                                     backColor=HexColor("#1A1D27"),
                                     textColor=HexColor("#E2E8F0"),
                                     leading=12, leftIndent=12, spaceAfter=6)

        story = [
            Paragraph("Fly OS Export", title_style),
            Paragraph(f"Generated by Fly OS v5  ·  {TODAY}", meta_style),
            HRFlowable(width="100%", thickness=0.5, color=HexColor("#2D3748")),
            Spacer(1, 10),
        ]

        in_code = False
        code_lines = []
        for line in text.split("\n"):
            if line.startswith("```"):
                if in_code:
                    story.append(Preformatted("\n".join(code_lines), code_style))
                    code_lines = []
                    in_code = False
                else:
                    in_code = True
                continue
            if in_code:
                code_lines.append(line)
                continue
            if line.startswith("# "):
                story.append(Paragraph(line[2:], h1_style))
            elif line.startswith("## "):
                story.append(Paragraph(line[3:], h2_style))
            elif line.startswith("### "):
                story.append(Paragraph(line[4:], h3_style))
            elif line.strip():
                safe = (line.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;"))
                story.append(Paragraph(safe, body_style))
            else:
                story.append(Spacer(1, 6))

        doc.build(story)
        return path
    except ImportError:
        raise RuntimeError("reportlab not installed.")


def export_pptx(text: str, filename: str) -> str:
    """Convert markdown-ish text into a PowerPoint presentation."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]  # blank
        title_layout = prs.slide_layouts[0]  # title slide

        PURPLE = RGBColor(0x6C, 0x63, 0xFF)
        WHITE  = RGBColor(0xF1, 0xF5, 0xF9)
        GRAY   = RGBColor(0x94, 0xA3, 0xB8)
        DARK   = RGBColor(0x0F, 0x11, 0x17)

        def add_title_slide(title_text, subtitle_text=""):
            slide = prs.slides.add_slide(title_layout)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = DARK
            t = slide.shapes.title
            t.text = title_text
            t.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
            t.text_frame.paragraphs[0].runs[0].font.size = Pt(40)
            if slide.placeholders[1] and subtitle_text:
                slide.placeholders[1].text = subtitle_text
                slide.placeholders[1].text_frame.paragraphs[0].runs[0].font.color.rgb = GRAY

        def add_content_slide(title_text, bullets):
            slide = prs.slides.add_slide(blank_layout)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = DARK

            # Title bar
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.9))
            tf = title_box.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = title_text
            p.runs[0].font.size  = Pt(28)
            p.runs[0].font.bold  = True
            p.runs[0].font.color.rgb = PURPLE

            # Accent line
            line = slide.shapes.add_shape(
                1, Inches(0.5), Inches(1.15), Inches(12), Emu(36000)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = PURPLE
            line.line.fill.background()

            # Bullets
            content_box = slide.shapes.add_textbox(
                Inches(0.6), Inches(1.4), Inches(12), Inches(5.7)
            )
            ctf = content_box.text_frame
            ctf.word_wrap = True
            first = True
            for bullet in bullets:
                if first:
                    bp = ctf.paragraphs[0]
                    first = False
                else:
                    bp = ctf.add_paragraph()
                bullet_text = bullet.lstrip("-•* ").strip()
                if bullet_text:
                    bp.text  = "  " + bullet_text
                    bp.level = 1 if bullet.startswith(("  ", "\t")) else 0
                    run = bp.runs[0]
                    run.font.size  = Pt(18)
                    run.font.color.rgb = WHITE
                    bp.space_before = Pt(6)

        # Parse text into slides
        slides_data = []
        current_title   = "Fly OS Presentation"
        current_bullets = []
        first_slide = True

        for line in text.split("\n"):
            if line.startswith("# ") or line.startswith("## "):
                if not first_slide and current_bullets:
                    slides_data.append((current_title, current_bullets))
                elif first_slide:
                    first_slide = False
                current_title   = line.lstrip("#").strip()
                current_bullets = []
            elif line.strip().startswith(("-", "•", "*")) or line.strip():
                if line.strip():
                    current_bullets.append(line)

        if current_bullets or current_title:
            slides_data.append((current_title, current_bullets))

        if not slides_data:
            slides_data = [("Fly OS Export", [text])]

        # Build first slide as title slide
        if slides_data:
            first_title   = slides_data[0][0]
            first_bullets = slides_data[0][1]
            # Subtitle: join first 1-2 bullets into a short line, or use date
            if first_bullets:
                subtitle = " · ".join(b.lstrip("-•* ").strip() for b in first_bullets[:2])
                subtitle = subtitle[:120]
            else:
                subtitle = f"Generated {TODAY}"
            add_title_slide(first_title, subtitle)
        for title, bullets in slides_data[1:]:
            add_content_slide(title, bullets)

        path = os.path.join(WORKSPACE_DIR, filename)
        prs.save(path)
        return path
    except ImportError:
        raise RuntimeError("python-pptx not installed.")


def export_xlsx(text: str, filename: str) -> str:
    """Convert tabular text/CSV into an Excel workbook."""
    try:
        import openpyxl
        from openpyxl.styles import PatternFill, Font, Alignment, Border, Side
        from openpyxl.utils import get_column_letter

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Fly OS Export"

        HEADER_FILL = PatternFill("solid", fgColor="6C63FF")
        HEADER_FONT = Font(bold=True, color="FFFFFF", name="Calibri", size=11)
        ALT_FILL    = PatternFill("solid", fgColor="1A1D27")
        BODY_FONT   = Font(color="F1F5F9", name="Calibri", size=11)
        BORDER_SIDE = Side(style="thin", color="2D3748")
        THIN_BORDER = Border(
            left=BORDER_SIDE, right=BORDER_SIDE,
            top=BORDER_SIDE, bottom=BORDER_SIDE
        )

        lines = [l for l in text.strip().split("\n") if l.strip()]
        row_idx = 1
        for i, line in enumerate(lines):
            # Detect CSV or pipe-delimited rows
            if "|" in line:
                cells = [c.strip() for c in line.split("|") if c.strip()]
            elif "," in line:
                import csv, io
                cells = next(csv.reader(io.StringIO(line)))
            else:
                cells = [line]

            for j, val in enumerate(cells, 1):
                cell = ws.cell(row=row_idx, column=j, value=val)
                cell.border = THIN_BORDER
                cell.alignment = Alignment(wrap_text=True, vertical="top")
                if i == 0:
                    cell.fill = HEADER_FILL
                    cell.font = HEADER_FONT
                else:
                    cell.font = BODY_FONT
                    if i % 2 == 0:
                        cell.fill = ALT_FILL

            # Skip markdown table divider rows
            if set(line.replace("|","").replace("-","").replace(":","").strip()) == set():
                ws.delete_rows(row_idx)
                continue
            row_idx += 1

        # Auto-size columns
        for col in ws.columns:
            max_len = max((len(str(c.value or "")) for c in col), default=0)
            ws.column_dimensions[get_column_letter(col[0].column)].width = min(max_len + 4, 60)

        # Freeze header row
        ws.freeze_panes = "A2"

        path = os.path.join(WORKSPACE_DIR, filename)
        wb.save(path)
        return path
    except ImportError:
        raise RuntimeError("openpyxl not installed.")


# ── Memory ──────────────────────────────────────────────────────────────────────

class MemoryDB:
    def __init__(self):
        self.conn = sqlite3.connect(DB_PATH, check_same_thread=False)
        self._init_db()

    def _init_db(self):
        self.conn.executescript("""
            CREATE TABLE IF NOT EXISTS conversations (
                id        INTEGER PRIMARY KEY AUTOINCREMENT,
                session   TEXT NOT NULL,
                role      TEXT NOT NULL,
                content   TEXT,
                timestamp DATETIME DEFAULT CURRENT_TIMESTAMP
            );
            CREATE TABLE IF NOT EXISTS sessions (
                id           TEXT PRIMARY KEY,
                title        TEXT,
                mode         TEXT,
                engine       TEXT DEFAULT 'gemini',
                resp_mode    TEXT DEFAULT 'balanced',
                created      DATETIME DEFAULT CURRENT_TIMESTAMP,
                updated      DATETIME DEFAULT CURRENT_TIMESTAMP
            );
        """)
        self.conn.commit()

    def new_session(self, title: str, mode: str) -> str:
        sid = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
        self.conn.execute(
            "INSERT INTO sessions (id, title, mode) VALUES (?, ?, ?)",
            (sid, title[:80], mode)
        )
        self.conn.commit()
        return sid

    def save_turn(self, session_id: str, role: str, content: str):
        self.conn.execute(
            "INSERT INTO conversations (session, role, content) VALUES (?, ?, ?)",
            (session_id, role, content)
        )
        clean = content[:80].replace("\n", " ").strip()
        self.conn.execute(
            "UPDATE sessions SET updated=CURRENT_TIMESTAMP, title=? WHERE id=?",
            (clean, session_id)
        )
        self.conn.commit()

    def update_engine(self, session_id: str, engine: str):
        self.conn.execute("UPDATE sessions SET engine=? WHERE id=?", (engine, session_id))
        self.conn.commit()

    def get_sessions(self, limit: int = 60) -> list[dict]:
        rows = self.conn.execute(
            "SELECT id, title, mode, engine, updated FROM sessions ORDER BY updated DESC LIMIT ?",
            (limit,)
        ).fetchall()
        return [{"id": r[0], "title": r[1] or "Untitled", "mode": r[2],
                 "engine": r[3], "updated": r[4]} for r in rows]

    def get_history(self, session_id: str) -> list[dict]:
        rows = self.conn.execute(
            "SELECT role, content FROM conversations WHERE session=? ORDER BY id",
            (session_id,)
        ).fetchall()
        return [{"role": r[0], "content": r[1]} for r in rows]

    def delete_session(self, session_id: str):
        self.conn.execute("DELETE FROM conversations WHERE session=?", (session_id,))
        self.conn.execute("DELETE FROM sessions WHERE id=?", (session_id,))
        self.conn.commit()

    def search_sessions(self, query: str) -> list[dict]:
        q = f"%{query}%"
        rows = self.conn.execute(
            "SELECT DISTINCT s.id, s.title, s.mode FROM sessions s "
            "JOIN conversations c ON c.session=s.id "
            "WHERE c.content LIKE ? OR s.title LIKE ? ORDER BY s.updated DESC LIMIT 20",
            (q, q)
        ).fetchall()
        return [{"id": r[0], "title": r[1] or "Untitled", "mode": r[2]} for r in rows]


# ── Main client ─────────────────────────────────────────────────────────────────

class FlyClient:
    """Triple Hybrid Engine AI client — v5.0"""

    def __init__(self):
        self.gemini = genai.Client(
            api_key=GEMINI_API_KEY,
            http_options={"api_version": "v1beta"},
        )
        self.db            = MemoryDB()
        self.current_mode  = "chat"
        self.response_mode = "balanced"
        self.session_id    = None
        self.history       = []
        self.last_response = ""
        self.active_engine = "gemini"
        self._stop_event   = threading.Event()
        self._current_thread: Optional[threading.Thread] = None

    # ── Session ─────────────────────────────────────────────────────────────

    def new_session(self):
        self.session_id    = None
        self.history       = []
        self.last_response = ""
        self.active_engine = "gemini"

    def load_session(self, session_id: str):
        self.session_id = session_id
        self.history    = self.db.get_history(session_id)

    def set_mode(self, mode: str):
        self.current_mode = mode

    def set_response_mode(self, mode: str):
        self.response_mode = mode

    def stop(self):
        self._stop_event.set()

    # ── Public stream entry ─────────────────────────────────────────────────

    def stream_response(
        self,
        user_text:        str,
        attachments:      list[str],
        on_token:         Callable[[str], None],
        on_done:          Callable[[str], None],
        on_error:         Callable[[str], None],
        on_file_saved:    Optional[Callable[[str, str], None]] = None,
        on_engine_change: Optional[Callable[[str], None]] = None,
        use_search:       bool = False,
        use_code:         bool = False,
    ):
        self._stop_event.clear()
        t = threading.Thread(
            target=self._worker,
            args=(user_text, attachments, on_token, on_done, on_error,
                  on_file_saved, on_engine_change, use_search, use_code),
            daemon=True,
        )
        self._current_thread = t
        t.start()

    # ── Worker ──────────────────────────────────────────────────────────────

    def _worker(self, user_text, attachments, on_token, on_done, on_error,
                on_file_saved, on_engine_change, use_search, use_code):

        eff_text = user_text.strip() or (
            "Please analyze the attached file(s)." if attachments else "Hello"
        )

        # ── Detect URLs and fetch content (in worker thread — non-blocking) ──
        urls = extract_urls(eff_text)
        url_contexts = []
        for url in urls[:3]:
            on_token(f"[📡 Reading: {url[:60]}{'...' if len(url)>60 else ''}]\n")
            content = fetch_url_content(url, timeout=12)
            # Trim aggressively so we don't blow the context window
            if len(content) > 6000:
                content = content[:6000] + "\n... [truncated]"
            url_contexts.append(f"\n### Webpage content — {url}\n{content}\n")
        if url_contexts:
            eff_text = eff_text + "\n\n" + "\n".join(url_contexts)

        # Bail early if user stopped while URLs were being fetched
        if self._stop_event.is_set():
            on_done("")
            return

        # ── Engine 1: Gemini ───────────────────────────────────────────────
        try:
            full = self._gemini_stream(
                eff_text, attachments, on_token, on_file_saved, use_search, use_code
            )
            self.active_engine = "gemini"
            self._finalize(user_text, full, "gemini", on_done)
            return
        except Exception as e1:
            if self._stop_event.is_set():
                on_done("")
                return
            gemini_err = str(e1)

        # ── Engine 2: OpenRouter ───────────────────────────────────────────
        try:
            if on_engine_change:
                on_engine_change("openrouter")
            on_token(f"\n[⚡ Gemini unavailable — switching to OpenRouter]\n\n")
            full = self._openrouter_call(eff_text)
            on_token(full)
            self.active_engine = "openrouter"
            self._finalize(user_text, full, "openrouter", on_done)
            return
        except Exception as e2:
            if self._stop_event.is_set():
                on_done("")
                return
            openrouter_err = str(e2)

        # ── Engine 3: Base44 ───────────────────────────────────────────────
        try:
            if on_engine_change:
                on_engine_change("base44")
            on_token(f"\n[⚡ OpenRouter unavailable — switching to Base44]\n\n")
            full = self._base44_call(eff_text)
            on_token(full)
            self.active_engine = "base44"
            self._finalize(user_text, full, "base44", on_done)
            return
        except Exception as e3:
            on_error(
                f"All engines failed.\n\n"
                f"Gemini: {_short(gemini_err)}\n"
                f"OpenRouter: {_short(openrouter_err)}\n"
                f"Base44: {_short(str(e3))}"
            )

    # ── Gemini streaming ────────────────────────────────────────────────────

    def _gemini_stream(self, user_text, attachments, on_token, on_file_saved,
                       use_search, use_code) -> str:
        parts = []

        for path_str in attachments:
            if self._stop_event.is_set():
                return ""
            path   = Path(path_str)
            suffix = path.suffix.lower()
            try:
                if suffix in IMAGE_EXTENSIONS:
                    raw  = path.read_bytes()
                    mime, _ = mimetypes.guess_type(str(path))
                    parts.append(types.Part.from_bytes(data=raw, mime_type=mime or "image/jpeg"))

                elif suffix in VIDEO_EXTENSIONS:
                    # Extract frames for vision analysis
                    frames = extract_video_frames(str(path))
                    if frames:
                        parts.append(types.Part.from_text(
                            text=f"### Video: {path.name} — analyzing {len(frames)} frames"
                        ))
                        for fp in frames:
                            raw = Path(fp).read_bytes()
                            parts.append(types.Part.from_bytes(data=raw, mime_type="image/png"))
                    else:
                        parts.append(types.Part.from_text(
                            text=f"[Video file attached: {path.name} — install opencv-python for frame analysis]"
                        ))

                elif suffix == ".pdf":
                    raw = path.read_bytes()
                    parts.append(types.Part.from_bytes(data=raw, mime_type="application/pdf"))

                elif suffix in {".docx", ".doc"}:
                    try:
                        from docx import Document as DocxDocument
                        doc = DocxDocument(str(path))
                        txt = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                        parts.append(types.Part.from_text(
                            text=f"### File: {path.name}\n```\n{txt[:12000]}\n```"
                        ))
                    except Exception:
                        txt = path.read_text(encoding="utf-8", errors="replace")
                        parts.append(types.Part.from_text(
                            text=f"### File: {path.name}\n```\n{txt[:12000]}\n```"
                        ))

                elif suffix in {".xlsx", ".xls"}:
                    try:
                        import pandas as pd
                        df = pd.read_excel(str(path), nrows=200)
                        parts.append(types.Part.from_text(
                            text=f"### Spreadsheet: {path.name}\n{df.to_markdown(index=False)}"
                        ))
                    except Exception:
                        parts.append(types.Part.from_text(
                            text=f"[Could not read {path.name}]"
                        ))

                elif suffix in {".pptx", ".ppt"}:
                    try:
                        from pptx import Presentation
                        prs = Presentation(str(path))
                        slides_text = []
                        for i, slide in enumerate(prs.slides, 1):
                            texts = [s.text for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
                            if texts:
                                slides_text.append(f"Slide {i}: " + " | ".join(texts))
                        content = "\n".join(slides_text[:50])
                        parts.append(types.Part.from_text(
                            text=f"### Presentation: {path.name}\n{content}"
                        ))
                    except Exception:
                        parts.append(types.Part.from_text(
                            text=f"[Could not read {path.name}]"
                        ))

                elif suffix in DOCUMENT_EXTENSIONS:
                    txt = path.read_text(encoding="utf-8", errors="replace")
                    parts.append(types.Part.from_text(
                        text=f"### File: {path.name}\n```\n{txt[:15000]}\n```"
                    ))

                elif suffix in AUDIO_EXTENSIONS:
                    raw  = path.read_bytes()
                    mime = "audio/mpeg" if suffix == ".mp3" else f"audio/{suffix[1:]}"
                    parts.append(types.Part.from_bytes(data=raw, mime_type=mime))

            except Exception as e:
                parts.append(types.Part.from_text(text=f"[Cannot read {path.name}: {e}]"))

        parts.append(types.Part.from_text(text=user_text))

        # Build history
        contents = []
        for turn in self.history[-30:]:
            role = "model" if turn["role"] == "assistant" else "user"
            contents.append(types.Content(
                role=role,
                parts=[types.Part.from_text(text=turn["content"])]
            ))
        contents.append(types.Content(role="user", parts=parts))

        # Tools
        tools = []
        if use_search:
            tools.append(types.Tool(google_search=types.GoogleSearch()))
        if use_code:
            tools.append(types.Tool(code_execution=types.ToolCodeExecution()))

        # Response mode settings
        rm = RESPONSE_MODES.get(self.response_mode, RESPONSE_MODES["balanced"])

        # System prompt
        sys_prompt = MODE_PROMPTS.get(self.current_mode, MODE_PROMPTS["chat"])
        ws = os.path.abspath(WORKSPACE_DIR).replace("\\", "/")
        sys_prompt += (
            f"\n\nCRITICAL: Today's real date is {TODAY}. "
            "Always use this as the current date.\n"
            f"Workspace folder: {ws}/\n"
            "When generating files (DOCX, PDF, PPTX, XLSX, images, code), "
            "tell the user what was created and where. "
            "For URLs pasted in messages: the content has already been fetched and included — use it directly."
        )

        config = types.GenerateContentConfig(
            system_instruction=sys_prompt,
            max_output_tokens=rm["max_tokens"],
            temperature=rm["temperature"],
            tools=tools if tools else None,
        )

        full = ""
        stream = self.gemini.models.generate_content_stream(
            model=MODEL_GEMINI, contents=contents, config=config
        )
        for chunk in stream:
            if self._stop_event.is_set():
                break
            if not chunk.candidates:
                continue
            for part in chunk.candidates[0].content.parts:
                if part.text:
                    full += part.text
                    on_token(part.text)
                elif hasattr(part, "inline_data") and part.inline_data:
                    mime  = part.inline_data.mime_type or "application/octet-stream"
                    data  = part.inline_data.data
                    ext   = _mime_to_ext(mime)
                    fname = f"fly_gen_{datetime.now().strftime('%H%M%S')}{ext}"
                    fpath = os.path.join(WORKSPACE_DIR, fname)
                    raw_b = data if isinstance(data, bytes) else base64.b64decode(data)
                    with open(fpath, "wb") as f:
                        f.write(raw_b)
                    msg = f"\n\n[📁 File saved: {fname}]"
                    full += msg
                    on_token(msg)
                    if on_file_saved:
                        on_file_saved(fname, fpath)
        return full

    # ── OpenRouter ──────────────────────────────────────────────────────────

    def _openrouter_call(self, user_text: str) -> str:
        rm = RESPONSE_MODES.get(self.response_mode, RESPONSE_MODES["balanced"])
        messages = [{"role": "system", "content": MODE_PROMPTS[self.current_mode]
                     + f"\n\nToday is {TODAY}."}]
        for turn in self.history[-20:]:
            role = "assistant" if turn["role"] == "assistant" else "user"
            messages.append({"role": role, "content": turn["content"]})
        messages.append({"role": "user", "content": user_text})

        resp = requests.post(
            "https://openrouter.ai/api/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {OPENROUTER_API_KEY}",
                "HTTP-Referer":  "https://flyos.app",
                "X-Title":       "Fly OS",
                "Content-Type":  "application/json",
            },
            json={
                "model": MODEL_OPENROUTER,
                "messages": messages,
                "max_tokens": rm["max_tokens"],
                "temperature": rm["temperature"],
            },
            timeout=90,
        )
        data = resp.json()
        if resp.status_code == 200 and "choices" in data:
            return data["choices"][0]["message"]["content"]
        err = data.get("error", {}).get("message", resp.text[:200])
        raise Exception(f"OpenRouter {resp.status_code}: {err}")

    # ── Base44 ──────────────────────────────────────────────────────────────

    def _base44_call(self, user_text: str) -> str:
        url = (f"https://app.base44.com/api/agents/{BASE44_AGENT_ID}"
               f"/conversations/{BASE44_CONVO_ID}/messages")
        resp = requests.post(
            url,
            headers={"Content-Type": "application/json", "api_key": BASE44_API_KEY},
            json={"role": "user", "content": user_text or "Hello"},
            timeout=120,
        )
        if resp.status_code == 200:
            data = resp.json()
            return (data.get("content") or data.get("message") or str(data))
        raise Exception(f"Base44 {resp.status_code}: {resp.text[:200]}")

    # ── Image generation ────────────────────────────────────────────────────

    def generate_image(self, prompt: str,
                       on_done:  Callable[[str, str], None],
                       on_error: Callable[[str], None]):
        def _run():
            try:
                response = self.gemini.models.generate_images(
                    model="imagen-3.0-generate-002",
                    prompt=prompt,
                    config=types.GenerateImagesConfig(number_of_images=1),
                )
                if response.generated_images:
                    raw   = response.generated_images[0].image.image_bytes
                    fname = f"fly_image_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png"
                    fpath = os.path.join(WORKSPACE_DIR, fname)
                    with open(fpath, "wb") as f:
                        f.write(raw)
                    on_done(fname, fpath)
                else:
                    on_error("No image returned from Gemini Imagen.")
            except Exception as e:
                on_error(f"Image generation failed: {e}")
        threading.Thread(target=_run, daemon=True).start()

    # ── Export ──────────────────────────────────────────────────────────────

    def export_last_response(self, fmt: str,
                             on_done:  Callable[[str, str], None],
                             on_error: Callable[[str], None]):
        def _run():
            try:
                if not self.last_response:
                    on_error("Nothing to export yet — send a message first.")
                    return
                ts = datetime.now().strftime("%Y%m%d_%H%M%S")
                if fmt == "docx":
                    fpath = export_docx(self.last_response, f"fly_export_{ts}.docx")
                elif fmt == "pdf":
                    fpath = export_pdf(self.last_response, f"fly_export_{ts}.pdf")
                elif fmt == "pptx":
                    fpath = export_pptx(self.last_response, f"fly_export_{ts}.pptx")
                elif fmt == "xlsx":
                    fpath = export_xlsx(self.last_response, f"fly_export_{ts}.xlsx")
                else:
                    on_error(f"Unknown format: {fmt}")
                    return
                on_done(Path(fpath).name, fpath)
            except Exception as e:
                on_error(str(e))
        threading.Thread(target=_run, daemon=True).start()

    # ── Background code execution ────────────────────────────────────────────

    def run_code(self, code: str,
                 on_done:  Callable[[str], None],
                 on_error: Callable[[str], None]):
        def _run():
            try:
                with tempfile.NamedTemporaryFile(
                    mode="w", suffix=".py", delete=False, encoding="utf-8"
                ) as f:
                    f.write(code)
                    tmp_path = f.name
                result = subprocess.run(
                    [sys.executable, tmp_path],
                    capture_output=True, text=True, timeout=30
                )
                os.unlink(tmp_path)
                output = ""
                if result.stdout:
                    output += result.stdout
                if result.stderr:
                    output += "\n[stderr]\n" + result.stderr
                on_done(output.strip() or "(No output)")
            except subprocess.TimeoutExpired:
                on_error("Code execution timed out (30s limit).")
            except Exception as e:
                on_error(str(e))
        threading.Thread(target=_run, daemon=True).start()

    # ── Auto-update check ────────────────────────────────────────────────────

    def check_for_update(self) -> Optional[dict]:
        try:
            resp = requests.get(GITHUB_VERSION_URL, timeout=5)
            if resp.status_code == 200:
                data = resp.json()
                remote = tuple(int(x) for x in data["version"].split("."))
                local  = tuple(int(x) for x in APP_VERSION.split("."))
                if remote > local:
                    return data
        except Exception:
            pass
        return None

    # ── Finalize ─────────────────────────────────────────────────────────────

    def _finalize(self, user_text: str, ai_text: str, engine: str, on_done):
        self.last_response = ai_text
        if not self.session_id:
            self.session_id = self.db.new_session(user_text[:80], self.current_mode)
        self.db.save_turn(self.session_id, "user", user_text)
        self.db.save_turn(self.session_id, "assistant", ai_text)
        self.db.update_engine(self.session_id, engine)
        self.history.append({"role": "user",      "content": user_text})
        self.history.append({"role": "assistant", "content": ai_text})
        # Cap in-memory history to last 60 turns to prevent unbounded growth
        if len(self.history) > 60:
            self.history = self.history[-60:]
        on_done(ai_text)


# ── Helpers ────────────────────────────────────────────────────────────────────

def _short(msg: str, n: int = 80) -> str:
    return msg[:n] + "..." if len(msg) > n else msg

def _mime_to_ext(mime: str) -> str:
    for k, v in {
        "image/png": ".png", "image/jpeg": ".jpg", "image/gif": ".gif",
        "image/webp": ".webp", "application/pdf": ".pdf",
        "text/csv": ".csv", "text/plain": ".txt", "application/json": ".json",
        "audio/mpeg": ".mp3", "audio/wav": ".wav",
    }.items():
        if k in mime:
            return v
    return ".bin"
